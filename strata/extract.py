import csv
import email
import json
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from .config import REPO_ROOT


def _tool_candidates(env_name: str, bundled_relative: str, executable_name: str) -> list[Path]:
    candidates: list[Path] = []
    env_value = os.environ.get(env_name)
    if env_value:
        candidates.append(Path(env_value))
    candidates.append(REPO_ROOT / "tools" / bundled_relative)
    which = shutil.which(executable_name)
    if which:
        candidates.append(Path(which))
    return candidates


def resolve_tesseract() -> str | None:
    for candidate in _tool_candidates("STRATA_TESSERACT", "tesseract/tesseract.exe", "tesseract"):
        if candidate.exists():
            return str(candidate)
    return None


def resolve_poppler() -> str | None:
    for candidate in _tool_candidates("STRATA_POPPLER_BIN", "poppler/bin", "pdftoppm"):
        if candidate.exists():
            return str(candidate if candidate.is_dir() else candidate.parent)
    return None


def resolve_libreoffice() -> str | None:
    for candidate in _tool_candidates("STRATA_LIBREOFFICE", "libreoffice/program/soffice.exe", "soffice"):
        if candidate.exists():
            return str(candidate)
    return None


def chunk_text_blocks(text: str, prefix: str, page: int | None = None) -> list[dict]:
    blocks = []
    paragraphs = [part.strip() for part in text.replace("\r", "").split("\n\n") if part.strip()]
    if not paragraphs and text.strip():
        paragraphs = [text.strip()]
    for index, paragraph in enumerate(paragraphs, start=1):
        snippet = paragraph[:240].replace("\n", " ")
        label = f"{prefix}{index}"
        blocks.append({"id": label, "label": label, "page": page, "text": paragraph, "snippet": snippet})
    return blocks


def _run_ocr_on_image(image) -> str:
    import pytesseract

    tesseract_cmd = resolve_tesseract()
    if tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
    return pytesseract.image_to_string(image, lang="eng")


def extract_pdf(path: Path) -> list[dict]:
    import pdfplumber

    chunks: list[dict] = []
    with pdfplumber.open(path) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if len(text) < 80:
                try:
                    from pdf2image import convert_from_path

                    images = convert_from_path(
                        str(path),
                        dpi=200,
                        first_page=page_number,
                        last_page=page_number,
                        poppler_path=resolve_poppler(),
                    )
                    if images:
                        ocr_text = _run_ocr_on_image(images[0]).strip()
                        if len(ocr_text) > len(text):
                            text = ocr_text
                except Exception:
                    pass
            if text:
                chunks.extend(chunk_text_blocks(text, prefix=f"p{page_number}-", page=page_number))
    return chunks


def extract_docx(path: Path) -> list[dict]:
    from docx import Document

    document = Document(str(path))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return chunk_text_blocks("\n\n".join(parts), prefix="para")


def extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    chunks: list[dict] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            text = "\n\n".join(texts)
            chunks.extend(chunk_text_blocks(text, prefix=f"slide{slide_number}-", page=slide_number))
    return chunks


def extract_csv(path: Path) -> list[dict]:
    rows = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        for row_index, row in enumerate(csv.reader(handle), start=1):
            value = " ".join(str(cell) for cell in row if cell is not None).strip()
            if value:
                rows.append({"id": f"row{row_index}", "label": f"row{row_index}", "page": None, "text": value, "snippet": value[:240]})
    return rows


def extract_xlsx(path: Path) -> list[dict]:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[dict] = []
    try:
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if values:
                    text = " ".join(values)
                    label = f"{sheet_name}!row{row_index}"
                    chunks.append({"id": label, "label": label, "page": None, "text": text, "snippet": text[:240]})
    finally:
        workbook.close()
    return chunks


def extract_text(path: Path) -> list[dict]:
    return chunk_text_blocks(path.read_text(encoding="utf-8", errors="replace"), prefix="para")


def extract_eml(path: Path) -> list[dict]:
    message = email.message_from_bytes(path.read_bytes())
    parts = []
    for header in ("From", "To", "Cc", "Subject", "Date"):
        value = message.get(header)
        if value:
            parts.append(f"{header}: {value}")
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True) or b""
                parts.append(payload.decode("utf-8", errors="replace"))
    else:
        payload = message.get_payload(decode=True) or b""
        parts.append(payload.decode("utf-8", errors="replace"))
    return chunk_text_blocks("\n\n".join(parts), prefix="para")


def extract_msg(path: Path) -> list[dict]:
    import extract_msg

    message = extract_msg.Message(str(path))
    parts = [item for item in [message.sender, message.to, message.subject, message.body] if item]
    return chunk_text_blocks("\n\n".join(parts), prefix="para")


def extract_image(path: Path) -> list[dict]:
    from PIL import Image

    with Image.open(path) as image:
        text = _run_ocr_on_image(image).strip()
    return chunk_text_blocks(text, prefix="ocr") if text else []


def convert_with_libreoffice(path: Path, expected_suffix: str) -> Path | None:
    soffice = resolve_libreoffice()
    if not soffice:
        return None
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_output = Path(temp_dir)
        subprocess.run(
            [soffice, "--headless", "--convert-to", expected_suffix.lstrip("."), "--outdir", str(temp_output), str(path)],
            check=False,
            capture_output=True,
            text=True,
        )
        converted = temp_output / f"{path.stem}.{expected_suffix.lstrip('.')}"
        if converted.exists():
            final_path = Path(tempfile.mkstemp(suffix=f".{expected_suffix.lstrip('.')}", prefix="strata-convert-")[1])
            final_path.write_bytes(converted.read_bytes())
            return final_path
    return None


def extract_document(path: Path) -> dict:
    extension = path.suffix.lower()
    title = path.name
    if extension == ".pdf":
        chunks = extract_pdf(path)
    elif extension == ".docx":
        chunks = extract_docx(path)
    elif extension == ".pptx":
        chunks = extract_pptx(path)
    elif extension in {".csv"}:
        chunks = extract_csv(path)
    elif extension in {".xlsx"}:
        chunks = extract_xlsx(path)
    elif extension in {".txt", ".md", ".rtf"}:
        chunks = extract_text(path)
    elif extension == ".eml":
        chunks = extract_eml(path)
    elif extension == ".msg":
        chunks = extract_msg(path)
    elif extension in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}:
        chunks = extract_image(path)
    elif extension in {".doc", ".xls", ".ppt"}:
        target_suffix = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}[extension]
        converted = convert_with_libreoffice(path, target_suffix)
        if not converted:
            raise RuntimeError(f"LibreOffice conversion is required for {extension} files.")
        try:
            converted_payload = extract_document(converted)
            chunks = converted_payload["chunks"]
        finally:
            converted.unlink(missing_ok=True)
    else:
        chunks = []
    return {"title": title, "chunks": chunks}


def save_index(index_path: Path, payload: dict) -> None:
    index_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")


def load_index(index_path: Path) -> dict:
    return json.loads(index_path.read_text(encoding="utf-8"))
