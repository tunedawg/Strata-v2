"""
Strata — Flask web server
Serves the UI on localhost and handles all API calls via JSON endpoints.
Run via run_windows.bat or run_mac.command — do not run directly.
"""

import os, sys, re, json, pickle, threading, time, difflib
import email, shutil, zipfile, io, csv, webbrowser
from pathlib import Path
from flask import Flask, request, jsonify, send_from_directory, send_file

# ── CONFIG ────────────────────────────────────────────────────────────────────
BASE_DIR     = os.path.join(os.path.expanduser("~"), "Documents", "Strata")
DATASETS_DIR = os.path.join(BASE_DIR, "datasets")
EXPORTS_DIR  = os.path.join(BASE_DIR, "exports")
TEMPLATE_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")
PORT         = int(os.environ.get("FLASK_PORT", 5000))

os.makedirs(DATASETS_DIR, exist_ok=True)
os.makedirs(EXPORTS_DIR,  exist_ok=True)

# ── Flask app ─────────────────────────────────────────────────────────────────
app = Flask(__name__, template_folder=TEMPLATE_DIR, static_folder=TEMPLATE_DIR)
app.config["MAX_CONTENT_LENGTH"] = 500 * 1024 * 1024  # 500MB upload limit

# ── In-memory state ───────────────────────────────────────────────────────────
LOADED_DATASET  = None
INDEX           = {}
INDEX_READY     = False
INDEX_PROGRESS  = {"status": "idle", "done": 0, "total": 0, "dataset": ""}
INDEX_LOCK      = threading.Lock()
REDACT_CANCEL   = False
REDACT_PROGRESS = {"status": "idle", "done": 0, "total": 0, "current": ""}
REDACT_RESULTS  = []
REDACT_LOCK     = threading.Lock()

# ── Lazy imports for heavy packages ───────────────────────────────────────────
def import_pdfplumber():
    import pdfplumber; return pdfplumber

def import_pypdf():
    import pypdf; return pypdf

def import_docx():
    from docx import Document; return Document

def import_openpyxl():
    import openpyxl; return openpyxl

def import_pptx():
    from pptx import Presentation; return Presentation

def import_reportlab():
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.units import inch
    return letter, landscape, SimpleDocTemplate, Paragraph, Spacer, HRFlowable, Table, TableStyle, getSampleStyleSheet, ParagraphStyle, colors, inch

# ── OCR ───────────────────────────────────────────────────────────────────────
OCR_ENABLED = True
OCR_MIN_CHARS = 100
_ocr_ready = None

def ocr_available():
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False

def run_ocr(img):
    try:
        import pytesseract
        return pytesseract.image_to_string(img, lang="eng")
    except Exception:
        return ""

def pdf_to_image(path, page_num):
    try:
        from pdf2image import convert_from_path
        imgs = convert_from_path(path, dpi=200, first_page=page_num+1, last_page=page_num+1)
        return imgs[0] if imgs else None
    except Exception:
        return None

# ── EXTRACTORS ────────────────────────────────────────────────────────────────
def extract_pdf(path):
    global _ocr_ready
    if _ocr_ready is None: _ocr_ready = ocr_available() and OCR_ENABLED
    chunks = {}
    try:
        pdfplumber = import_pdfplumber()
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = (page.extract_text() or "").strip()
                if len(text) < OCR_MIN_CHARS and _ocr_ready:
                    img = pdf_to_image(path, i)
                    if img:
                        ocr = run_ocr(img).strip()
                        if len(ocr) > len(text): text = ocr
                if text.strip():
                    chunks[f"p{i+1}"] = text.lower()
    except Exception as e:
        print(f"  PDF error {path}: {e}")
    return chunks

def extract_docx(path):
    chunks = {}
    try:
        Document = import_docx()
        doc = Document(path)
        buf = []; idx = 1
        for para in doc.paragraphs:
            t = para.text.strip()
            if t: buf.append(t)
            if len(buf) >= 5:
                chunks[f"para{idx}"] = " ".join(buf).lower(); idx += 1; buf = []
        for table in doc.tables:
            for row in table.rows:
                rt = " ".join(c.text.strip() for c in row.cells if c.text.strip())
                if rt: buf.append(rt)
                if len(buf) >= 5:
                    chunks[f"para{idx}"] = " ".join(buf).lower(); idx += 1; buf = []
        if buf: chunks[f"para{idx}"] = " ".join(buf).lower()
    except Exception as e:
        print(f"  DOCX error {path}: {e}")
    return chunks

def extract_pptx(path):
    chunks = {}
    try:
        Presentation = import_pptx()
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts: chunks[f"slide{i+1}"] = " ".join(texts).lower()
    except Exception as e:
        print(f"  PPTX error {path}: {e}")
    return chunks

def extract_csv(path):
    chunks = {}
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            for i, row in enumerate(csv.reader(f)):
                text = " ".join(str(c) for c in row).strip()
                if text: chunks[f"row{i+1}"] = text.lower()
    except Exception as e:
        print(f"  CSV error {path}: {e}")
    return chunks

def extract_xlsx(path):
    chunks = {}
    try:
        openpyxl = import_openpyxl()
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        try:
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    parts = [str(c).strip() for c in row if c is not None]
                    text = " ".join(p for p in parts if p)
                    if text: chunks[f"{sheet}_row{i+1}"] = text.lower()
        finally:
            wb.close()
    except Exception as e:
        print(f"  XLSX error {path}: {e}")
    return chunks

def extract_txt(path):
    chunks = {}
    try:
        text = Path(path).read_text(encoding="utf-8", errors="replace")
        for i, p in enumerate(re.split(r"\n\s*\n", text)):
            if p.strip(): chunks[f"para{i+1}"] = p.strip().lower()
    except Exception as e:
        print(f"  TXT error {path}: {e}")
    return chunks

def extract_eml(path):
    chunks = {}
    try:
        with open(path, "rb") as f:
            msg = email.message_from_bytes(f.read())
        parts = []
        for h in ("From","To","Cc","Subject","Date"):
            v = msg.get(h, "")
            if v: parts.append(f"{h}: {v}")
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    try: parts.append(part.get_payload(decode=True).decode("utf-8", errors="replace"))
                    except: pass
        else:
            try: parts.append(msg.get_payload(decode=True).decode("utf-8", errors="replace"))
            except: pass
        full = "\n".join(parts)
        for i, p in enumerate(re.split(r"\n\s*\n", full)):
            if p.strip(): chunks[f"para{i+1}"] = p.strip().lower()
    except Exception as e:
        print(f"  EML error {path}: {e}")
    return chunks

def extract_msg(path):
    chunks = {}
    try:
        import extract_msg as em
        msg = em.Message(path)
        parts = []
        if msg.sender:  parts.append(f"From: {msg.sender}")
        if msg.to:      parts.append(f"To: {msg.to}")
        if msg.subject: parts.append(f"Subject: {msg.subject}")
        if msg.body:    parts.append(msg.body)
        full = "\n".join(parts)
        for i, p in enumerate(re.split(r"\n\s*\n", full)):
            if p.strip(): chunks[f"para{i+1}"] = p.strip().lower()
        msg.close()
    except Exception as e:
        print(f"  MSG error {path}: {e}")
    return chunks

def extract_image(path):
    global _ocr_ready
    if _ocr_ready is None: _ocr_ready = ocr_available() and OCR_ENABLED
    chunks = {}
    if not _ocr_ready: return chunks
    try:
        from PIL import Image
        text = run_ocr(Image.open(path)).strip()
        if text: chunks["p1"] = text.lower()
    except Exception as e:
        print(f"  Image OCR error {path}: {e}")
    return chunks

EXTRACTORS = {
    ".pdf": extract_pdf, ".docx": extract_docx, ".doc": extract_docx,
    ".pptx": extract_pptx, ".ppt": extract_pptx,
    ".csv": extract_csv, ".xlsx": extract_xlsx, ".xls": extract_xlsx,
    ".txt": extract_txt, ".md": extract_txt, ".rtf": extract_txt,
    ".eml": extract_eml, ".msg": extract_msg,
    ".png": extract_image, ".jpg": extract_image, ".jpeg": extract_image,
    ".tif": extract_image, ".tiff": extract_image, ".bmp": extract_image,
}

# ── INDEXING ──────────────────────────────────────────────────────────────────
def index_path(dataset_name):
    return os.path.join(DATASETS_DIR, dataset_name, "_index.pkl")

def get_datasets():
    datasets = []
    for d in sorted(Path(DATASETS_DIR).iterdir()):
        if d.is_dir():
            idx = index_path(d.name)
            datasets.append({
                "name": d.name,
                "indexed": os.path.exists(idx),
                "files": len([f for f in d.rglob("*")
                              if f.is_file() and not f.name.startswith("_")
                              and f.suffix.lower() in EXTRACTORS]),
            })
    return datasets

def build_index_for(dataset_name):
    global INDEX, INDEX_READY, INDEX_PROGRESS, LOADED_DATASET
    folder = Path(DATASETS_DIR) / dataset_name
    files  = [f for f in folder.rglob("*")
              if f.is_file() and not f.name.startswith("_")
              and f.suffix.lower() in EXTRACTORS]
    with INDEX_LOCK:
        INDEX_PROGRESS = {"status": "indexing", "done": 0, "total": len(files), "dataset": dataset_name}
        INDEX_READY = False
    local_index = {}
    for i, fpath in enumerate(files):
        try:
            chunks = EXTRACTORS[fpath.suffix.lower()](str(fpath))
            rel = str(fpath.relative_to(folder))
            local_index[rel] = chunks
            with INDEX_LOCK:
                INDEX_PROGRESS["done"] = i + 1
        except Exception as e:
            print(f"  ERROR on {fpath.name}: {e}")
    with open(index_path(dataset_name), "wb") as f:
        pickle.dump(local_index, f)
    with INDEX_LOCK:
        INDEX = local_index; INDEX_READY = True
        LOADED_DATASET = dataset_name; INDEX_PROGRESS["status"] = "ready"

def load_index_for(dataset_name):
    global INDEX, INDEX_READY, INDEX_PROGRESS, LOADED_DATASET
    with INDEX_LOCK:
        INDEX_PROGRESS = {"status": "loading", "done": 0, "total": 0, "dataset": dataset_name}
        INDEX_READY = False
    with open(index_path(dataset_name), "rb") as f:
        local_index = pickle.load(f)
    with INDEX_LOCK:
        INDEX = local_index; INDEX_READY = True; LOADED_DATASET = dataset_name
        INDEX_PROGRESS = {"status": "ready", "done": len(local_index),
                          "total": len(local_index), "dataset": dataset_name}

# ── QUERY ENGINE ──────────────────────────────────────────────────────────────
def words_of(text):
    return re.findall(r"[a-z0-9']+", text.lower())

def match_term(text, token):
    token = token.strip()
    if token.startswith('"') and token.endswith('"'):
        return token[1:-1].lower() in text
    fuzzy = token.startswith("~") or token.endswith("~")
    token = token.strip("~")
    if "*" in token:
        pattern = re.compile(r"\b" + re.escape(token).replace(r"\*", r"\w*") + r"\b")
        return bool(pattern.search(text))
    if fuzzy:
        return bool(difflib.get_close_matches(token.lower(), words_of(text), n=1, cutoff=0.75))
    return token.lower() in text

def proximity_match(text, phrase_a, phrase_b, within):
    def find_pos(wlist, phrase):
        pw = words_of(phrase)
        n = len(pw)
        return [i for i in range(len(wlist)-n+1) if wlist[i:i+n] == pw]
    wlist = words_of(text)
    pos_a = find_pos(wlist, phrase_a.strip('"'))
    pos_b = find_pos(wlist, phrase_b.strip('"'))
    if not pos_a or not pos_b: return False
    return any(abs(a-b) <= within for a in pos_a for b in pos_b)

def parse_query(query):
    prox_re = re.compile(r'(".*?"|[\w~*]+)\s+(?:W|NEAR)/(\d+)\s+(".*?"|[\w~*]+)', re.IGNORECASE)
    def split_top(q, op):
        parts, depth, in_quote, buf = [], 0, False, ""
        i = 0
        while i < len(q):
            c = q[i]
            if c == '"': in_quote = not in_quote; buf += c
            elif not in_quote and c == "(": depth += 1; buf += c
            elif not in_quote and c == ")": depth -= 1; buf += c
            elif not in_quote and depth == 0 and q[i:i+len(op)].upper() == op.upper():
                parts.append(buf.strip()); buf = ""; i += len(op); continue
            else: buf += c
            i += 1
        parts.append(buf.strip())
        return [p for p in parts if p]
    def matching_close(s):
        depth = 0
        for i, c in enumerate(s):
            if c == "(": depth += 1
            elif c == ")":
                depth -= 1
                if depth == 0: return i
        return -1
    def evaluate(text, q):
        q = q.strip()
        if not q: return False
        parts = split_top(q, " OR ")
        if len(parts) > 1: return any(evaluate(text, p) for p in parts)
        parts = split_top(q, " AND NOT ")
        if len(parts) > 1: return evaluate(text, parts[0]) and not any(evaluate(text, p) for p in parts[1:])
        parts = split_top(q, " AND ")
        if len(parts) > 1: return all(evaluate(text, p) for p in parts)
        parts = split_top(q, " NOT ")
        if len(parts) > 1: return evaluate(text, parts[0]) and not any(evaluate(text, p) for p in parts[1:])
        if q.upper().startswith("NOT "): return not evaluate(text, q[4:])
        if q.startswith("(") and q.endswith(")") and matching_close(q) == len(q)-1:
            return evaluate(text, q[1:-1])
        pm = prox_re.search(q)
        if pm: return proximity_match(text, pm.group(1), pm.group(3), int(pm.group(2)))
        return match_term(text, q)
    return lambda text: evaluate(text, query)

def run_search(terms, index):
    results = {}
    for term in terms:
        term = term.strip()
        if not term: continue
        try:
            matcher = parse_query(term)
        except Exception as e:
            results[term] = {"error": str(e), "total_hits": 0, "docs": []}
            continue
        docs_hit = []
        total = 0
        for doc_name, chunks in index.items():
            matched = [cid for cid, text in chunks.items() if matcher(text)]
            if matched:
                docs_hit.append({"name": doc_name, "chunks": matched})
                total += len(matched)
        docs_hit.sort(key=lambda x: -len(x["chunks"]))
        results[term] = {"total_hits": total, "docs": docs_hit}
    return results

# ── REDACTION DETECTION ───────────────────────────────────────────────────────
def _color_is_black(c):
    if c is None: return False
    try:
        if isinstance(c, (int, float)): return float(c) < 0.15
        if isinstance(c, (list, tuple)):
            if len(c) == 1: return float(c[0]) < 0.15
            if len(c) == 3: return all(float(v) < 0.15 for v in c)
            if len(c) == 4: return float(c[3]) > 0.85
    except: pass
    return False

def detect_redactions_in_pdf(path):
    findings = []
    try:
        import pypdf
        reader = pypdf.PdfReader(path, strict=False)
        for page_num, page in enumerate(reader.pages, 1):
            annots = page.get("/Annots")
            if not annots: continue
            for ref in annots:
                try:
                    a = ref.get_object() if hasattr(ref, "get_object") else ref
                    if str(a.get("/Subtype","")) == "/Redact":
                        rect = a.get("/Rect", [])
                        findings.append({
                            "page": page_num, "type": "PDF Redaction Annotation",
                            "location": f"Page {page_num}",
                            "details": "Standard /Redact annotation",
                        })
                except: pass
    except: pass
    try:
        pdfplumber = import_pdfplumber()
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                pw = float(page.width or 612); ph = float(page.height or 792)
                for rect in (page.rects or []):
                    try:
                        x0  = float(rect.get("x0",0)); x1 = float(rect.get("x1",0))
                        top = float(rect.get("top",0)); bot = float(rect.get("bottom",top))
                        w = abs(x1-x0); h = abs(bot-top)
                        if w < 15 or h < 3: continue
                        fill = rect.get("non_stroking_color")
                        if fill is None: continue
                        if _color_is_black(fill) and h <= 72:
                            if h < 4 or (w/pw > 0.95 and h/ph > 0.95): continue
                            findings.append({
                                "page": page_num, "type": "Black Redaction Bar",
                                "location": f"Page {page_num} — x:{x0:.0f}–{x1:.0f}, y:{top:.0f}–{bot:.0f}",
                                "details": f"{w:.0f}×{h:.0f}pt black rectangle",
                            })
                    except: pass
    except: pass
    black_locs = {f["location"] for f in findings if f["type"] == "Black Redaction Bar"}
    findings = [f for f in findings if not (f["type"] == "PDF Redaction Annotation" and f["location"] in black_locs)]
    return findings

def run_redaction_scan(dataset_name):
    global REDACT_PROGRESS, REDACT_RESULTS, REDACT_CANCEL
    folder = Path(DATASETS_DIR) / dataset_name
    files  = sorted([f for f in folder.rglob("*")
                     if f.is_file() and not f.name.startswith("_")
                     and f.suffix.lower() in {".pdf", ".docx", ".doc"}])
    with REDACT_LOCK:
        REDACT_PROGRESS = {"status":"scanning","done":0,"total":len(files),"current":"","dataset":dataset_name}
        REDACT_RESULTS  = []; REDACT_CANCEL = False
    results = []
    for i, fpath in enumerate(files):
        with REDACT_LOCK:
            if REDACT_CANCEL:
                REDACT_PROGRESS["status"] = "cancelled"; REDACT_RESULTS = results; return
            REDACT_PROGRESS["done"] = i; REDACT_PROGRESS["current"] = fpath.name
        try:
            findings = detect_redactions_in_pdf(str(fpath))
            if findings:
                pages = sorted({f["page"] for f in findings if f.get("page")})
                type_counts = {}
                for f in findings: type_counts[f["type"]] = type_counts.get(f["type"],0) + 1
                results.append({
                    "filename": str(fpath.relative_to(folder)),
                    "ext": fpath.suffix.lower(), "redaction_count": len(findings),
                    "pages": pages, "type_summary": type_counts, "findings": findings,
                })
            with REDACT_LOCK: REDACT_RESULTS = results.copy()
        except Exception as e:
            print(f"  Scan error {fpath.name}: {e}")
    with REDACT_LOCK:
        REDACT_PROGRESS["status"] = "done"; REDACT_PROGRESS["done"] = len(files)
        REDACT_RESULTS = results

# ── FLASK ROUTES ──────────────────────────────────────────────────────────────

# Favicon
@app.route("/favicon.ico")
def favicon():
    return send_from_directory(TEMPLATE_DIR, "favicon.ico", mimetype="image/x-icon")

@app.route("/favicon.png")
def favicon_png():
    return send_from_directory(TEMPLATE_DIR, "favicon.png", mimetype="image/png")

# Serve HTML pages
@app.route("/")
def index():
    return send_from_directory(TEMPLATE_DIR, "index.html")

@app.route("/production")
def production():
    return send_from_directory(TEMPLATE_DIR, "production.html")

@app.route("/redactions")
def redactions_page():
    return send_from_directory(TEMPLATE_DIR, "redactions.html")

# Preview a file
@app.route("/preview/<dataset>/<path:filename>")
def preview(dataset, filename):
    folder = str(Path(DATASETS_DIR) / dataset)
    return send_from_directory(folder, filename)

# ── Dataset routes ────────────────────────────────────────────────────────────
@app.route("/datasets")
def datasets():
    return jsonify(get_datasets())

@app.route("/datasets/create", methods=["POST"])
def create_dataset():
    name = re.sub(r'[\\/*?:"<>|]', "_", (request.json.get("name") or "").strip())
    if not name: return jsonify({"error": "Name required"}), 400
    path = Path(DATASETS_DIR) / name
    if path.exists(): return jsonify({"error": "Already exists"}), 400
    path.mkdir(parents=True)
    return jsonify({"ok": True, "name": name})

@app.route("/datasets/delete", methods=["POST"])
def delete_dataset():
    global INDEX, INDEX_READY, LOADED_DATASET, INDEX_PROGRESS
    name = request.json.get("name")
    folder = Path(DATASETS_DIR) / name
    if not folder.exists(): return jsonify({"error": "Not found"}), 404
    shutil.rmtree(str(folder))
    if LOADED_DATASET == name:
        with INDEX_LOCK:
            INDEX = {}; INDEX_READY = False; LOADED_DATASET = None
            INDEX_PROGRESS = {"status":"idle","done":0,"total":0,"dataset":""}
    return jsonify({"ok": True})

@app.route("/upload", methods=["POST"])
def upload():
    dataset_name = request.form.get("dataset")
    folder = Path(DATASETS_DIR) / dataset_name
    if not folder.exists(): return jsonify({"error": "Dataset not found"}), 404
    saved = []
    for f in request.files.getlist("files"):
        dest = folder / f.filename
        dest.parent.mkdir(parents=True, exist_ok=True)
        f.save(str(dest))
        saved.append(f.filename)
    return jsonify({"ok": True, "saved": saved})

@app.route("/index/build", methods=["POST"])
def build_index():
    dataset_name = request.json.get("dataset")
    folder = Path(DATASETS_DIR) / dataset_name
    if not folder.exists(): return jsonify({"error": "Not found"}), 404
    t = threading.Thread(target=build_index_for, args=(dataset_name,), daemon=True)
    t.start()
    return jsonify({"ok": True})

@app.route("/index/load", methods=["POST"])
def load_index():
    dataset_name = request.json.get("dataset")
    if not os.path.exists(index_path(dataset_name)):
        return jsonify({"error": "Not indexed"}), 404
    t = threading.Thread(target=load_index_for, args=(dataset_name,), daemon=True)
    t.start()
    return jsonify({"ok": True})

@app.route("/index/status")
def index_status():
    with INDEX_LOCK:
        return jsonify({**INDEX_PROGRESS, "ready": INDEX_READY, "loaded": LOADED_DATASET})

# ── Search routes ─────────────────────────────────────────────────────────────
@app.route("/search", methods=["POST"])
def search():
    if not INDEX_READY: return jsonify({"error": "No dataset loaded"}), 400
    terms = request.json.get("terms", [])
    results = run_search(terms, INDEX)
    return jsonify({"results": results, "dataset": LOADED_DATASET})

@app.route("/production/search", methods=["POST"])
def production_search():
    data = request.json
    dataset_name = data.get("dataset")
    terms = data.get("terms", [])
    pkl = index_path(dataset_name)
    if not os.path.exists(pkl): return jsonify({"error": "Not indexed"}), 404
    with open(pkl, "rb") as f:
        local_index = pickle.load(f)
    results = run_search(terms, local_index)
    return jsonify({"dataset": dataset_name, "results": results})

# ── Export routes ─────────────────────────────────────────────────────────────
@app.route("/export/hit_report", methods=["POST"])
def export_hit_report():
    if not INDEX_READY: return jsonify({"error": "No dataset loaded"}), 400
    data    = request.json
    results = data.get("results", {})
    fmt     = data.get("fmt", "xlsx")
    from datetime import datetime as dt
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
    path = os.path.join(EXPORTS_DIR, f"hit_report_{timestamp}.xlsx")
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Hit Report"
    hdr_font = Font(bold=True, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="1A1A2E")
    stripe   = PatternFill("solid", fgColor="F0F4FF")
    for col, h in enumerate(["Search Term","Total Hits","Documents","Pages"], 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = hdr_font; cell.fill = hdr_fill
        cell.alignment = Alignment(horizontal="center")
    ws.column_dimensions["A"].width = 40
    ws.column_dimensions["B"].width = 14
    ws.column_dimensions["C"].width = 14
    ws.column_dimensions["D"].width = 50
    for row_idx, (term, r) in enumerate(results.items(), 2):
        docs = r.get("docs", [])
        pages = []
        for d in docs:
            for c in d.get("chunks", []):
                if re.match(r'^p\d+$', c): pages.append(f"{d['name']}:{c}")
        ws.cell(row=row_idx, column=1, value=term)
        ws.cell(row=row_idx, column=2, value=r.get("total_hits",0)).alignment = Alignment(horizontal="center")
        ws.cell(row=row_idx, column=3, value=len(docs)).alignment = Alignment(horizontal="center")
        ws.cell(row=row_idx, column=4, value=", ".join(pages[:20]))
        if row_idx % 2 == 0:
            for col in range(1,5): ws.cell(row=row_idx, column=col).fill = stripe
    wb.save(path)
    return send_file(path, as_attachment=True, download_name=f"hit_report_{timestamp}.xlsx")

@app.route("/production/export_zip", methods=["POST"])
def export_production_zip():
    data = request.json
    dataset_name = data.get("dataset")
    selections   = data.get("selections", [])
    folder = Path(DATASETS_DIR) / dataset_name
    if not folder.exists(): return jsonify({"error": "Dataset not found"}), 404
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from datetime import datetime as dt
    timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Responsiveness Log"
    hdr_font = Font(bold=True, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="1A1A2E")
    for col, h in enumerate(["Filename","Responsive Terms","Pages"], 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = hdr_font; cell.fill = hdr_fill
    ws.column_dimensions["A"].width = 50
    ws.column_dimensions["B"].width = 55
    ws.column_dimensions["C"].width = 35
    for row_idx, sel in enumerate(selections, 2):
        ws.cell(row=row_idx, column=1, value=sel.get("name",""))
        ws.cell(row=row_idx, column=2, value=", ".join(sel.get("terms",[])))
        ws.cell(row=row_idx, column=3, value=", ".join(sel.get("chunks",[])))
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        xlsx_buf = io.BytesIO(); wb.save(xlsx_buf); xlsx_buf.seek(0)
        zf.writestr("responsiveness_log.xlsx", xlsx_buf.read())
        for sel in selections:
            rel_path = sel.get("name","")
            src = folder / rel_path
            if src.exists():
                zf.write(str(src), arcname="files/" + rel_path.replace("\\","/"))
    zip_buf.seek(0)
    return send_file(zip_buf, mimetype="application/zip", as_attachment=True,
                     download_name=f"production_export_{timestamp}.zip")

# ── Redaction routes ──────────────────────────────────────────────────────────
@app.route("/redactions/scan", methods=["POST"])
def start_scan():
    dataset_name = request.json.get("dataset")
    if not (Path(DATASETS_DIR) / dataset_name).exists():
        return jsonify({"error": "Dataset not found"}), 404
    t = threading.Thread(target=run_redaction_scan, args=(dataset_name,), daemon=True)
    t.start()
    return jsonify({"status": "started"})

@app.route("/redactions/status")
def scan_status():
    with REDACT_LOCK:
        prog = {**REDACT_PROGRESS, "results_count": len(REDACT_RESULTS)}
    done = prog.get("done", 0); total = prog.get("total", 1) or 1
    prog["percent"] = int(done / total * 100)
    return jsonify(prog)

@app.route("/redactions/results")
def scan_results():
    with REDACT_LOCK:
        results = list(REDACT_RESULTS)
    return jsonify({"documents": results,
                    "total_findings": sum(r.get("redaction_count",0) for r in results)})

@app.route("/redactions/cancel", methods=["POST"])
def cancel_scan():
    global REDACT_CANCEL
    with REDACT_LOCK:
        REDACT_CANCEL = True
    return jsonify({"ok": True})

@app.route("/redactions/export_zip", methods=["POST"])
def export_redaction_zip():
    with REDACT_LOCK:
        ds      = REDACT_PROGRESS.get("dataset","")
        results = list(REDACT_RESULTS)
    folder = Path(DATASETS_DIR) / ds
    from datetime import datetime as dt
    timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for r in results:
            src = folder / r["filename"]
            if src.exists():
                zf.write(str(src), arcname="redacted_docs/" + r["filename"].replace("\\","/"))
    zip_buf.seek(0)
    return send_file(zip_buf, mimetype="application/zip", as_attachment=True,
                     download_name=f"redacted_docs_{timestamp}.zip")

@app.route("/redactions/export_list", methods=["POST"])
def export_redaction_list():
    with REDACT_LOCK:
        results = list(REDACT_RESULTS)
        dataset = REDACT_PROGRESS.get("dataset","export")
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from datetime import datetime as dt
    timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Findings"
    hdr_font = Font(bold=True, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="1A1A2E")
    stripe   = PatternFill("solid", fgColor="FFF0F0")
    for col, (h,w) in enumerate([("Document",45),("Page",8),("Type",28),("Location",38),("Details",55)],1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = hdr_font; cell.fill = hdr_fill
        ws.column_dimensions[cell.column_letter].width = w
    row_idx = 2
    for r in results:
        for f in r.get("findings",[]):
            ws.cell(row=row_idx, column=1, value=r["filename"])
            ws.cell(row=row_idx, column=2, value=f.get("page","—"))
            ws.cell(row=row_idx, column=3, value=f["type"])
            ws.cell(row=row_idx, column=4, value=f["location"])
            ws.cell(row=row_idx, column=5, value=f["details"])
            if row_idx % 2 == 0:
                for col in range(1,6): ws.cell(row=row_idx, column=col).fill = stripe
            row_idx += 1
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                     as_attachment=True, download_name=f"redaction_findings_{timestamp}.xlsx")

@app.route("/redactions/export_report", methods=["POST"])
def export_redaction_report():
    fmt = request.args.get("format", "xlsx")
    with REDACT_LOCK:
        results = list(REDACT_RESULTS)
        dataset = REDACT_PROGRESS.get("dataset","export")
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from datetime import datetime as dt
    timestamp = dt.now().strftime("%Y%m%d_%H%M%S")
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Summary"
    hdr_font = Font(bold=True, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="1A1A2E")
    stripe   = PatternFill("solid", fgColor="FFF0F0")
    for col, (h,w) in enumerate([("Document",52),("Count",18),("Pages",30),("Types",55)],1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = hdr_font; cell.fill = hdr_fill
        ws.column_dimensions[cell.column_letter].width = w
    for row_idx, r in enumerate(results, 2):
        pages_str = ", ".join(f"p.{p}" for p in r["pages"][:20])
        types_str = "; ".join(f"{t}({c})" for t,c in r["type_summary"].items())
        ws.cell(row=row_idx, column=1, value=r["filename"])
        ws.cell(row=row_idx, column=2, value=r["redaction_count"])
        ws.cell(row=row_idx, column=3, value=pages_str)
        ws.cell(row=row_idx, column=4, value=types_str)
        if row_idx % 2 == 0:
            for col in range(1,5): ws.cell(row=row_idx, column=col).fill = stripe
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    return send_file(buf, mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                     as_attachment=True, download_name=f"redaction_report_{timestamp}.xlsx")

# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=" * 50)
    print("  Strata — Legal Document Search")
    print("=" * 50)
    print(f"  Data folder : {BASE_DIR}")
    print(f"  URL         : http://127.0.0.1:{PORT}")
    print("=" * 50)
    print()
    app.run(host="127.0.0.1", port=PORT, debug=False, threaded=True)
